"""
╔══════════════════════════════════════════════════════════════════════════════╗
║  SCHNEIDER ELECTRIC - RULE EXTRACTOR v3.3                                   ║
║  Uses raw HTTP requests — works regardless of SDK versions installed        ║
║  Priority: OpenAI → Gemini → Ollama                                        ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""

import streamlit as st
import json
import os
import subprocess
import requests
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv

# ── File parsers ──────────────────────────────────────────────────────────────
try:
    import pdfplumber
    PDF_OK = True
except ImportError:
    PDF_OK = False

try:
    from pptx import Presentation
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    from docx import Document
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

# ── Config ────────────────────────────────────────────────────────────────────
load_dotenv()

RULES_GUIDE_FOLDER = Path("Rules_Guide_Used")
RULES_OUTPUT_FILE  = Path("server/Extracted_Rules_From_Pdf.json")
GEMINI_API_KEY     = os.getenv("GEMINI_API_KEY", "")
OPENAI_API_KEY     = os.getenv("OPENAI_API_KEY", "")
OLLAMA_BASE_URL    = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
OLLAMA_MODEL       = os.getenv("OLLAMA_MODEL", "codellama")
GITHUB_REPO_PATH   = Path(".")

RULES_GUIDE_FOLDER.mkdir(exist_ok=True)

# ── Page setup ────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Rule Extractor v3.3 | Schneider AI",
    page_icon="📋",
    layout="wide"
)

st.markdown("""
<style>
    .stButton>button { background:#3DCD58; color:white; border:none; border-radius:8px; font-weight:bold; }
    .stButton>button:hover { background:#2db347; }
    .rule-card { background:#1e2130; border-left:4px solid #3DCD58; border-radius:8px; padding:12px 16px; margin:8px 0; }
    .file-card { background:#1a1d2e; border:1px solid #2d3250; border-radius:8px; padding:10px 14px; margin:6px 0; }
    .fallback-badge { display:inline-block; padding:4px 10px; border-radius:12px; font-size:11px; font-weight:bold; margin-left:8px; }
    .badge-gemini { background:#4285F4; color:white; }
    .badge-openai { background:#10A37F; color:white; }
    .badge-ollama { background:#FF6F61; color:white; }
    h1 { color:#3DCD58 !important; }
</style>
""", unsafe_allow_html=True)

# ══════════════════════════════════════════════════════════════════════════════
# RAW HTTP LLM CLIENT — no SDK dependencies, always works
# ══════════════════════════════════════════════════════════════════════════════

class TripleLLM:
    """
    Calls AI APIs via raw HTTP requests.
    Zero dependency on openai or google-generativeai SDK versions.
    Priority: OpenAI → Gemini → Ollama
    """

    def __init__(self):
        self.last_used = None
        self.init_errors = []

        # ── OpenAI check ─────────────────────────────────────────────────────
        self.openai_ready = False
        if OPENAI_API_KEY:
            try:
                r = requests.get(
                    "https://api.openai.com/v1/models",
                    headers={"Authorization": f"Bearer {OPENAI_API_KEY}"},
                    timeout=8
                )
                if r.status_code == 200:
                    self.openai_ready = True
                    st.success("✅ OpenAI ready (PRIMARY) — connected via API")
                elif r.status_code == 401:
                    msg = "OpenAI API key is invalid"
                    self.init_errors.append(msg)
                    st.warning(f"⚠️ {msg}")
                else:
                    msg = f"OpenAI returned HTTP {r.status_code}"
                    self.init_errors.append(msg)
                    st.warning(f"⚠️ {msg}")
            except Exception as e:
                msg = f"OpenAI connection failed: {str(e)}"
                self.init_errors.append(msg)
                st.warning(f"⚠️ {msg}")
        else:
            st.info("ℹ️ No OPENAI_API_KEY — add it to Streamlit Secrets")

        # ── Gemini check ─────────────────────────────────────────────────────
        self.gemini_ready = False
        self.gemini_model = None
        if GEMINI_API_KEY:
            # Try each model with a real ping
            for model in ["gemini-1.5-flash", "gemini-1.5-pro", "gemini-pro"]:
                try:
                    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={GEMINI_API_KEY}"
                    r = requests.post(
                        url,
                        json={"contents": [{"parts": [{"text": "hi"}]}]},
                        timeout=10
                    )
                    if r.status_code == 200:
                        self.gemini_ready = True
                        self.gemini_model = model
                        st.success(f"✅ Gemini ready ({model})")
                        break
                    elif r.status_code == 400:
                        # Model exists but bad request format — still usable
                        self.gemini_ready = True
                        self.gemini_model = model
                        st.success(f"✅ Gemini ready ({model})")
                        break
                except Exception:
                    continue
            if not self.gemini_ready:
                msg = "Gemini: no working model found (check API key)"
                self.init_errors.append(msg)
                st.warning(f"⚠️ {msg}")
        else:
            st.info("ℹ️ No GEMINI_API_KEY — add it to Streamlit Secrets")

        # ── Ollama check ─────────────────────────────────────────────────────
        self.ollama_ready = False
        try:
            r = requests.get(f"{OLLAMA_BASE_URL}/api/tags", timeout=3)
            if r.status_code == 200:
                models = r.json().get("models", [])
                names = [m.get("name","").split(":")[0] for m in models]
                if any(OLLAMA_MODEL in n for n in names):
                    self.ollama_ready = True
                    st.success(f"✅ Ollama ready — model: {OLLAMA_MODEL}")
                else:
                    st.warning(f"⚠️ Ollama running but '{OLLAMA_MODEL}' not found. Run: ollama pull {OLLAMA_MODEL}")
        except Exception:
            st.info("ℹ️ Ollama not running (optional)")

    def _call_openai(self, prompt: str) -> str:
        """Call OpenAI via raw HTTP — works with ANY openai SDK version installed"""
        response = requests.post(
            "https://api.openai.com/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENAI_API_KEY}",
                "Content-Type": "application/json"
            },
            json={
                "model": "gpt-4o-mini",
                "messages": [
                    {"role": "system", "content": "You are a Schneider Electric coding standards expert. Extract coding rules. Return ONLY valid JSON arrays."},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.3,
                "max_tokens": 4000
            },
            timeout=60
        )
        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"].strip()

    def _call_gemini(self, prompt: str) -> str:
        """Call Gemini via raw HTTP"""
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{self.gemini_model}:generateContent?key={GEMINI_API_KEY}"
        response = requests.post(
            url,
            json={"contents": [{"parts": [{"text": prompt}]}]},
            timeout=60
        )
        response.raise_for_status()
        data = response.json()
        return data["candidates"][0]["content"]["parts"][0]["text"].strip()

    def _call_ollama(self, prompt: str) -> str:
        """Call Ollama via raw HTTP"""
        response = requests.post(
            f"{OLLAMA_BASE_URL}/api/generate",
            json={
                "model": OLLAMA_MODEL,
                "prompt": prompt,
                "stream": False,
                "options": {"temperature": 0.3, "num_predict": 2000}
            },
            timeout=120
        )
        response.raise_for_status()
        return response.json().get("response", "")

    def generate(self, prompt: str) -> tuple:
        """
        Try OpenAI → Gemini → Ollama.
        Returns: (success, text, model_used)
        """
        # 1. OpenAI
        if self.openai_ready:
            try:
                text = self._call_openai(prompt)
                self.last_used = "openai"
                return True, text, "openai"
            except Exception as e:
                err = str(e).lower()
                if "quota" in err or "rate_limit" in err:
                    st.warning("⚠️ OpenAI quota exceeded, trying Gemini...")
                elif "401" in err:
                    st.error("❌ OpenAI key invalid")
                else:
                    st.warning(f"⚠️ OpenAI failed: {str(e)[:100]}, trying Gemini...")

        # 2. Gemini
        if self.gemini_ready:
            try:
                text = self._call_gemini(prompt)
                self.last_used = "gemini"
                return True, text, "gemini"
            except Exception as e:
                err = str(e).lower()
                if "quota" in err or "429" in err:
                    st.warning("⚠️ Gemini quota exceeded, trying Ollama...")
                else:
                    st.warning(f"⚠️ Gemini failed: {str(e)[:100]}, trying Ollama...")

        # 3. Ollama
        if self.ollama_ready:
            try:
                with st.spinner(f"🤖 Using Ollama ({OLLAMA_MODEL})... 30-60 seconds..."):
                    text = self._call_ollama(prompt)
                self.last_used = "ollama"
                return True, text, "ollama"
            except requests.exceptions.Timeout:
                st.error("❌ Ollama timeout. Try smaller model.")
            except Exception as e:
                st.error(f"❌ Ollama failed: {str(e)}")

        st.error("❌ All AI providers failed!")
        return False, "", "none"

    def get_status(self) -> dict:
        return {
            "openai": self.openai_ready,
            "gemini": self.gemini_ready,
            "ollama": self.ollama_ready,
            "last_used": self.last_used,
            "errors": self.init_errors
        }

@st.cache_resource
def get_triple_llm():
    return TripleLLM()

# ── Text extractors ───────────────────────────────────────────────────────────
def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    try:
        if ext == ".pdf":
            if not PDF_OK:
                return "ERROR: run pip install pdfplumber"
            text = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        text.append(t)
            return "\n".join(text)
        elif ext in (".pptx", ".ppt"):
            if not PPTX_OK:
                return "ERROR: run pip install python-pptx"
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text.strip())
            return "\n".join(text)
        elif ext in (".docx", ".doc"):
            if not DOCX_OK:
                return "ERROR: run pip install python-docx"
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == ".txt":
            return open(file_path, encoding="utf-8", errors="ignore").read()
        else:
            return f"ERROR: Unsupported type {ext}"
    except Exception as e:
        return f"ERROR: {e}"

# ── AI extraction ─────────────────────────────────────────────────────────────
def extract_rules_with_ai(text: str, source_file: str, llm: TripleLLM) -> tuple:
    prompt = f"""Extract coding rules from this document.

For each rule, return a JSON object with:
- "rule_id": "NEW_001"
- "rule": The coding rule (1-2 sentences)
- "suggested_fix": How to comply
- "source": "{source_file}"
- "category": One of [naming, structure, security, energy, documentation, safety, performance, general]
- "severity": One of [critical, error, warning, info]

Return ONLY a JSON array. No markdown, no explanation. If no rules found, return [].

Document text (first 8000 chars):
{text[:8000]}"""

    success, raw, model_used = llm.generate(prompt)
    if not success:
        return [], "none"

    try:
        cleaned = raw.strip()
        if "```" in cleaned:
            lines = cleaned.split("\n")
            json_lines, in_block = [], False
            for line in lines:
                if line.startswith("```"):
                    in_block = not in_block
                    continue
                json_lines.append(line)
            cleaned = "\n".join(json_lines).strip()

        rules = json.loads(cleaned)
        if isinstance(rules, list):
            return rules, model_used
        st.warning(f"⚠️ Response not a list: {type(rules)}")
        return [], model_used
    except json.JSONDecodeError as e:
        st.error(f"❌ JSON parse error: {str(e)}")
        with st.expander("🔍 Raw response"):
            st.code(raw[:1000])
        return [], model_used

# ── Rules JSON helpers ────────────────────────────────────────────────────────
def load_existing_rules() -> list:
    if RULES_OUTPUT_FILE.exists():
        try:
            return json.loads(RULES_OUTPUT_FILE.read_text(encoding="utf-8")).get("rules", [])
        except Exception:
            return []
    return []

def save_rules(rules: list):
    RULES_OUTPUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    RULES_OUTPUT_FILE.write_text(
        json.dumps({"rules": rules}, indent=2, ensure_ascii=False),
        encoding="utf-8"
    )

def get_next_rule_number(existing: list) -> int:
    nums = ["".join(filter(str.isdigit, r.get("rule_id", ""))) for r in existing]
    nums = [int(n) for n in nums if n]
    return max(nums, default=0) + 1

def renumber(rules: list, start: int) -> list:
    for i, r in enumerate(rules):
        r["rule_id"] = f"R{start + i:03d}"
    return rules

# ── Git helpers ───────────────────────────────────────────────────────────────
def run_git(cmd: list) -> tuple:
    try:
        r = subprocess.run(cmd, cwd=str(GITHUB_REPO_PATH.resolve()),
                           capture_output=True, text=True, timeout=30)
        return r.returncode == 0, r.stdout + r.stderr
    except Exception as e:
        return False, str(e)

def push_to_github(commit_msg: str) -> tuple:
    steps = []
    ok, out = run_git(["git", "add", str(RULES_OUTPUT_FILE)])
    steps.append(("git add", ok, out))
    if not ok:
        return False, steps
    ok, out = run_git(["git", "commit", "-m", commit_msg])
    steps.append(("git commit", ok, out))
    if not ok:
        if "nothing to commit" in out:
            steps.append(("note", True, "Already up to date."))
            return True, steps
        return False, steps
    ok, out = run_git(["git", "push", "origin", "main"])
    steps.append(("git push", ok, out))
    return ok, steps

def scan_folder() -> list:
    exts = {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".txt"}
    return sorted([f for f in RULES_GUIDE_FOLDER.iterdir()
                   if f.is_file() and f.suffix.lower() in exts])

# ═══════════════════════════════════════════════════════════════════════════════
# UI
# ═══════════════════════════════════════════════════════════════════════════════
st.title("📋 Rule Extractor v3.3")
st.markdown("**Triple-LLM Fallback:** OpenAI (Primary) → Gemini → Ollama")
st.divider()

llm = get_triple_llm()
status = llm.get_status()

st.subheader("🤖 AI Provider Status")

if status.get("errors"):
    with st.expander("⚠️ Initialization Warnings", expanded=True):
        for error in status["errors"]:
            st.warning(error)

col_o, col_g, col_ol = st.columns(3)

with col_o:
    if status["openai"]:
        st.success("✅ **OpenAI** Ready (PRIMARY)")
    else:
        st.warning("⚠️ **OpenAI** Unavailable")
        st.caption("Add OPENAI_API_KEY to Streamlit Secrets")

with col_g:
    if status["gemini"]:
        st.success("✅ **Gemini** Ready")
    else:
        st.warning("⚠️ **Gemini** Unavailable")
        st.caption("Add GEMINI_API_KEY to Streamlit Secrets")

with col_ol:
    if status["ollama"]:
        st.success("✅ **Ollama** Ready")
        st.caption(f"Model: {OLLAMA_MODEL}")
    else:
        st.warning("⚠️ **Ollama** Offline")
        with st.expander("How to fix"):
            st.code(f"ollama pull {OLLAMA_MODEL}\nollama serve")

if not any([status["openai"], status["gemini"], status["ollama"]]):
    st.error("❌ No AI providers available! Add API keys to Streamlit Secrets.")
    st.stop()

existing_rules = load_existing_rules()

st.divider()
c1, c2, c3, c4 = st.columns(4)
c1.metric("📋 Rules in JSON", len(existing_rules))
c2.metric("📁 Files in Folder", len(scan_folder()))
c3.metric("🆕 This Session", len(st.session_state.get("extracted_rules", [])))
if status["last_used"]:
    badge_class = f"badge-{status['last_used']}"
    c4.markdown(f'Last used: <span class="fallback-badge {badge_class}">{status["last_used"].upper()}</span>',
                unsafe_allow_html=True)
else:
    c4.metric("🤖 Status", "Ready")

st.divider()

col1, col2 = st.columns(2, gap="large")

with col1:
    st.subheader("📁 Upload Files")
    uploaded = st.file_uploader(
        "Drop PDF / PPT / Word / TXT files here",
        type=["pdf", "pptx", "ppt", "docx", "doc", "txt"],
        accept_multiple_files=True,
        label_visibility="collapsed"
    )
    if uploaded:
        for f in uploaded:
            dest = RULES_GUIDE_FOLDER / f.name
            dest.write_bytes(f.read())
            st.success(f"✅ Saved: **{f.name}**")

    st.divider()
    st.subheader("📂 Rules_Guide_Used")
    files = scan_folder()
    if not files:
        st.info("No files yet. Upload above.")
    else:
        for fp in files:
            ext = fp.suffix.upper().lstrip(".")
            icon = {"PDF":"📄","PPTX":"📊","PPT":"📊","DOCX":"📝","DOC":"📝","TXT":"📃"}.get(ext,"📄")
            kb = max(1, fp.stat().st_size // 1024)
            st.markdown(
                f'<div class="file-card">{icon} <b>{fp.name}</b> '
                f'<span style="color:#888;float:right">{kb} KB</span></div>',
                unsafe_allow_html=True
            )

with col2:
    st.subheader("🤖 Extract with AI")
    files = scan_folder()

    if not files:
        st.info("Upload files on the left first.")
    else:
        file_names = [f.name for f in files]
        selected = st.multiselect("Select files:", options=file_names, default=file_names)

        ca, cb = st.columns(2)
        auto_save  = ca.toggle("Auto-save to JSON", value=True)
        skip_dupes = cb.toggle("Skip duplicates",   value=True)

        if st.button("🚀 Extract Rules Now", use_container_width=True):
            if not selected:
                st.warning("Select at least one file")
            else:
                all_new = []
                model_usage = {"gemini": 0, "openai": 0, "ollama": 0}
                prog = st.progress(0, text="Starting...")

                for idx, fname in enumerate(selected):
                    prog.progress(idx / len(selected), text=f"📖 {fname}")
                    with st.spinner(f"Reading {fname}..."):
                        text = extract_text(str(RULES_GUIDE_FOLDER / fname))
                        if text.startswith("ERROR"):
                            st.warning(f"⚠️ {text}")
                            continue
                        rules, model_used = extract_rules_with_ai(text, fname, llm)
                        if rules:
                            all_new.extend(rules)
                            model_usage[model_used] = model_usage.get(model_used, 0) + 1
                            badge_class = f"badge-{model_used}"
                            st.markdown(
                                f'✅ **{fname}** → {len(rules)} rules '
                                f'<span class="fallback-badge {badge_class}">{model_used.upper()}</span>',
                                unsafe_allow_html=True
                            )
                        else:
                            st.warning(f"⚠️ No rules found in {fname}")

                prog.progress(1.0, text="✅ Done!")

                if any(model_usage.values()):
                    st.info(f"📊 Models used: Gemini: {model_usage['gemini']} | OpenAI: {model_usage['openai']} | Ollama: {model_usage['ollama']}")

                if all_new:
                    all_new = renumber(all_new, get_next_rule_number(existing_rules))
                    if skip_dupes:
                        exist_set = {r.get("rule","").lower() for r in existing_rules}
                        all_new = [r for r in all_new if r.get("rule","").lower() not in exist_set]

                    st.session_state["extracted_rules"] = all_new
                    st.info(f"📋 **{len(all_new)} unique new rules** ready")

                    if auto_save:
                        save_rules(existing_rules + all_new)
                        st.session_state["rules_saved"] = True
                        st.success(f"💾 Saved! Total: **{len(existing_rules)+len(all_new)}** rules")
                        st.balloons()
                else:
                    st.error("No rules extracted from any file.")

st.divider()
st.subheader("🐙 Push to GitHub")

extracted = st.session_state.get("extracted_rules", [])

if not extracted:
    st.info("Extract rules first, then push here.")
else:
    default_msg = (f"Add {len(extracted)} new rules from Rules_Guide_Used "
                   f"[{datetime.now().strftime('%Y-%m-%d %H:%M')}]")
    col_msg, col_btn = st.columns([3, 1])
    commit_msg = col_msg.text_input("Commit message:", value=default_msg, label_visibility="collapsed")

    if col_btn.button("⬆️ Push to GitHub", use_container_width=True):
        if not st.session_state.get("rules_saved"):
            save_rules(existing_rules + extracted)
        with st.spinner("Pushing..."):
            success, steps = push_to_github(commit_msg)
        for name, ok, out in steps:
            icon = "✅" if ok else "❌"
            with st.expander(f"{icon} {name}", expanded=not ok):
                st.code(out or "No output")
        if success:
            st.success("🎉 **Pushed to GitHub successfully!**")
            st.markdown("🔗 [View on GitHub](https://github.com/ShriHarsan64K/Schneider-AI-Code-Reviewer/blob/main/server/Extracted_Rules_From_Pdf.json)")
        else:
            st.error("❌ Push failed — check details above.")

st.divider()
st.subheader("👀 Preview Extracted Rules")

# ── Download buttons ──────────────────────────────────────────────────────────
dl1, dl2 = st.columns(2)

# Download extracted (this session)
if extracted:
    dl1.download_button(
        label="⬇️ Download Extracted Rules (this session)",
        data=json.dumps({"rules": extracted}, indent=2, ensure_ascii=False),
        file_name=f"extracted_rules_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
        mime="application/json",
        use_container_width=True
    )

# Download full JSON (all rules saved)
if existing_rules:
    dl2.download_button(
        label="⬇️ Download Full JSON (all rules)",
        data=json.dumps({"rules": existing_rules}, indent=2, ensure_ascii=False),
        file_name="Extracted_Rules_From_Pdf.json",
        mime="application/json",
        use_container_width=True
    )

if extracted:
    cats = ["All"] + sorted({r.get("category","general") for r in extracted})
    sel_cat = st.selectbox("Filter by category:", cats)
    filtered = extracted if sel_cat == "All" else [r for r in extracted if r.get("category") == sel_cat]
    st.markdown(f"Showing **{len(filtered)}** rules:")
    for rule in filtered[:50]:
        sev = rule.get("severity","info")
        col = {"critical":"#ff4444","error":"#ff8800","warning":"#ffcc00","info":"#3DCD58"}.get(sev,"#888")
        st.markdown(f"""
<div class="rule-card">
    <div style="display:flex;justify-content:space-between;margin-bottom:4px">
        <b style="color:#3DCD58">{rule.get('rule_id','')}</b>
        <span style="color:{col};font-size:12px;text-transform:uppercase">{sev}</span>
    </div>
    <div style="color:#e0e0e0;margin-bottom:6px">{rule.get('rule','')}</div>
    <div style="color:#888;font-size:13px">🔧 {rule.get('suggested_fix','')}</div>
    <div style="color:#555;font-size:11px;margin-top:4px">📁 {rule.get('source','')} · 🏷️ {rule.get('category','')}</div>
</div>""", unsafe_allow_html=True)
else:
    st.info("Run an extraction above to see rules here.")